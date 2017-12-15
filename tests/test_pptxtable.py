import pytest

from pptx_tables import PptxTable
from pptx.util import Inches


class TestPptxTable:
    data1 = [[0, 1, 2],
             [3, 4, 5],
             [6, 7, 8]]

    data2 = [{"apples": 0, "bananas": 1, "pears": 2},
             {"apples": 3, "bananas": 4, "pears": 5},
             {"apples": 6, "bananas": 7, "pears": 8}]

    def test_create_table_data1(self):
        tbl = PptxTable(self.data1)
        tbl.set_table_location(Inches(0), Inches(1), Inches(5), Inches(2))
        tbl.create_table(slide_index=0,
                         rows_sort_order=[2, 1, 0],
                         columns_sort_order=[2, 1, 0],
                         column_headers=["column2", "column1", "column0"])
        assert tbl.pptx_table.cell(0, 0).text_frame.text == "column2"
        assert tbl.pptx_table.cell(0, 1).text_frame.text == "column1"
        assert tbl.pptx_table.cell(0, 2).text_frame.text == "column0"
        assert tbl.pptx_table.cell(1, 0).text_frame.text == "8"
        assert tbl.pptx_table.cell(1, 1).text_frame.text == "7"
        assert tbl.pptx_table.cell(1, 2).text_frame.text == "6"
        assert tbl.pptx_table.cell(2, 0).text_frame.text == "5"
        assert tbl.pptx_table.cell(2, 1).text_frame.text == "4"
        assert tbl.pptx_table.cell(2, 2).text_frame.text == "3"
        assert tbl.pptx_table.cell(3, 0).text_frame.text == "2"
        assert tbl.pptx_table.cell(3, 1).text_frame.text == "1"
        assert tbl.pptx_table.cell(3, 2).text_frame.text == "0"
        for i in [0, 1, 2]:
            with pytest.raises(IndexError):
                print(tbl.pptx_table.cell(i, 3))
        for j in [0, 1, 2, 3]:
            with pytest.raises(IndexError):
                print(tbl.pptx_table.cell(4, j))

    def test_create_table_data2(self):
        tbl = PptxTable(self.data2)

        tbl.set_table_location(Inches(0), Inches(3), Inches(5), Inches(2))
        tbl.create_table(slide_index=0,
                         rows_sort_order=[2, 1, 0],
                         columns_sort_order=["pears", "bananas", "apples"],
                         column_headers=["Pears", "Bananas", "Apples"])
        assert tbl.pptx_table.cell(0, 0).text_frame.text == "Pears"
        assert tbl.pptx_table.cell(0, 1).text_frame.text == "Bananas"
        assert tbl.pptx_table.cell(0, 2).text_frame.text == "Apples"
        assert tbl.pptx_table.cell(1, 0).text_frame.text == "8"
        assert tbl.pptx_table.cell(1, 1).text_frame.text == "7"
        assert tbl.pptx_table.cell(1, 2).text_frame.text == "6"
        assert tbl.pptx_table.cell(2, 0).text_frame.text == "5"
        assert tbl.pptx_table.cell(2, 1).text_frame.text == "4"
        assert tbl.pptx_table.cell(2, 2).text_frame.text == "3"
        assert tbl.pptx_table.cell(3, 0).text_frame.text == "2"
        assert tbl.pptx_table.cell(3, 1).text_frame.text == "1"
        assert tbl.pptx_table.cell(3, 2).text_frame.text == "0"
        for i in [0, 1, 2]:
            with pytest.raises(IndexError):
                print(tbl.pptx_table.cell(i, 3))
        for j in [0, 1, 2, 3]:
            with pytest.raises(IndexError):
                print(tbl.pptx_table.cell(4, j))
