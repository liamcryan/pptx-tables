from pptx import Presentation
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt

from pptx_tables.collections import Collection


class PptxTable(object):
    """  This class represents a PowerPoint Table.

    Attributes:
        prs:  a Python-pptx presentation
        table_args:  the arguments needed for Python-pptx presentation.add_table method
        pptx_table:  the table returned from the Python-pptx presentation.add_table method
        collection:  the class representing data portion of the table

    """
    def __init__(self, data, presentation=None):
        """   Instantiate the class with data and an optional presentation.

        :param data: this is the data to be placed in the table
                    must look like [[1, 0], [0, 0], [2, 1]]
                    OR
                    must look like [{"item1": 1, "item2": 2}, {"item1": 3, "item2": 3}]
        :param presentation: this is the Python-pptx presentation.  If not provided, a default presentation
                            will be created.
        """
        self.prs = presentation
        self.table_args = [None, None, Inches(0), Inches(0), Inches(2), None]  # default placement
        self.pptx_table = None
        self.font_size = Pt(8)
        self.row_height = Inches(.38)
        self.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        self.first_row = True
        self.first_col = False

        self.collection = Collection(data=data)

    def _add_table(self, slide_index=0):
        """ Creates a slide if needed then adds table according to the table_args provided.

        :param slide_index: slide you want to add to...only really needed if you have an existing presentation
        :return: None
        """
        if not self.prs:
            if slide_index > 0:
                raise Exception("slide index provided is greater than the number of slides in the report")
            prs = Presentation()
            self.prs = prs
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        else:
            if slide_index > len(list(self.prs.slides._sldIdLst)) - 1:
                raise Exception("slide index provided is greater than the number of slides in the report")
            slide = self.prs.slides[slide_index]

        shapes = slide.shapes
        self.set_table_size(len(self.collection.rows.idx), len(self.collection.columns.idx))
        table = shapes.add_table(*self.table_args).table
        self.pptx_table = table

    def create_table(self, slide_index=0, rows_sort_order=None, columns_sort_order=None, columns_headers=None,
                     columns_widths_weight=None, transpose=False):
        """ Sorts the rows/columns. Provides column headers.  Creates table. Puts data in a table.

        :param slide_index:  what slide do you want to put the table on?
        :param rows_sort_order:  list, how to sort rows
        :param columns_sort_order: list, how to sort columns
        :param columns_headers: list, what to call columns (dependent upon sorting of columns) for example:
                                columns are sorted in this order : [0, 1, 2],
                                columns headers should be something like : ["column_0", "column_1", "column_2"]
                                columns are sorted in this order : [2, 1, 0],
                                columns headers should be something like : ["column_2", "column_1", "column_0"]
        :param columns_widths_weight:  list, what is the weight given to each column,
                                        sum of list should add to length of list to maintain table width
        :param transpose:  if True, the transpose will be displayed
        :return: None
        """
        if rows_sort_order:
            self.collection.rows.sort_order(rows_sort_order)
        if columns_sort_order:
            self.collection.columns.sort_order(columns_sort_order)
        if columns_headers:
            self.collection.set_column_headers(columns_headers)
        if transpose:
            self.collection.rows.idx, self.collection.columns.idx = \
                self.collection.columns.idx, self.collection.rows.idx

            self._add_table(slide_index)

            self.pptx_table.first_row = self.first_col
            self.pptx_table.first_col = self.first_row
        else:
            self._add_table(slide_index)
            self.pptx_table.first_row = self.first_row
            self.pptx_table.first_col = self.first_col

        if columns_widths_weight:
            self.set_columns_widths_weight(columns_widths_weight)

        for i, row in enumerate(self.collection.rows.idx):
            for j, col in enumerate(self.collection.columns.idx):
                if not transpose:
                    self.pptx_table.cell(i, j).text = str(self.collection.data[row][col])
                    self.pptx_table.cell(i, j).text_frame.paragraphs[0].font.size = self.font_size
                    self.pptx_table.cell(i, j).text_frame.paragraphs[0].alignment = self.alignment
                    self.pptx_table.cell(i, j).margin_top = 0
                    self.pptx_table.cell(i, j).vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
                else:
                    self.pptx_table.cell(i, j).text = str(self.collection.data[col][row])
                    self.pptx_table.cell(i, j).text_frame.paragraphs[0].font.size = self.font_size
                    self.pptx_table.cell(i, j).text_frame.paragraphs[0].alignment = self.alignment
                    self.pptx_table.cell(i, j).margin_top = 0
                    self.pptx_table.cell(i, j).vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

    def save_pptx(self, file_name):
        """ save a presentation """
        self.prs.save(file_name)

    def set_table_location(self, left, top, width, height=None):
        """ set the table's location """
        self.table_args = [None, None, left, top, width, height]

    def set_table_size(self, rows, columns):
        self.table_args[0] = rows
        self.table_args[1] = columns
        if not self.table_args[5]:
            height = self.row_height / 914400  # Inches(1) -> 914400
            self.table_args[5] = Inches(rows * height)

    def set_columns_widths_weight(self, column_widths):
        """ set the proportion of space each column takes """
        for j, col in enumerate(self.collection.columns.idx):
            table_width = self.table_args[4]/914400  # Inches(1) -> 914400
            table_columns = self.table_args[1]
            if isinstance(col, int):
                column_proportion = column_widths[col]
                self.pptx_table.columns[col].width = Inches(column_proportion * (table_width / table_columns))
            elif isinstance(col, str):
                column_proportion = column_widths[j]
                self.pptx_table.columns[j].width = Inches(column_proportion * (table_width / table_columns))

    def set_formatting(self, font_size=None, alignment=None, row_height=None, first_row=None, first_col=None):
        """ set the formatting, default is used if this method is not called """
        if font_size:
            self.font_size = font_size
        if alignment:
            self.alignment = alignment
        if row_height:
            self.row_height = row_height
        if first_row:
            self.first_row = first_row
        if first_col:
            self.first_col = first_col

if __name__ == "__main__":
    data1 = [[0, 1, 2], [3, 4, 5], [6, 7, 8]]

    tbl0 = PptxTable(data1)
    tbl0.set_table_location(left=Inches(0), top=Inches(1), width=Inches(5))
    tbl0.set_formatting(font_size=Pt(7), alignment=PP_PARAGRAPH_ALIGNMENT.LEFT, row_height=Inches(1))
    tbl0.create_table(slide_index=0,
                      rows_sort_order=[2, 1, 0],
                      columns_sort_order=[2, 1, 0],
                      columns_headers=["column2", "column1", "column0"],
                      columns_widths_weight=[.75, .75, 1.5], )
    tbl0.save_pptx("test0.pptx")

    tbl1 = PptxTable(data1)
    tbl1.set_table_location(left=Inches(0), top=Inches(1), width=Inches(5))
    tbl1.set_formatting(font_size=Pt(7), alignment=PP_PARAGRAPH_ALIGNMENT.LEFT, row_height=Inches(1))
    tbl1.create_table(slide_index=0,
                      rows_sort_order=[2, 1, 0],
                      columns_sort_order=[2, 1, 0],
                      columns_headers=["column2", "column1", "column0"],
                      columns_widths_weight=[1.6, .8, .8, .8],
                      transpose=True)
    tbl1.save_pptx("test1.pptx")

    data2 = [{"apples": 0, "bananas": 1, "pears": 2},
             {"apples": 3, "bananas": 4, "pears": 5},
             {"apples": 6, "bananas": 7, "pears": 8}]

    tbl2 = PptxTable(data2)

    tbl2.set_table_location(Inches(0), Inches(3), Inches(5), Inches(2))
    tbl2.create_table(slide_index=0,
                      rows_sort_order=[2, 1, 0],
                      columns_sort_order=["pears", "bananas", "apples"],
                      columns_headers=["Pears", "Bananas", "Apples"],
                      columns_widths_weight=[1, 1, 1])
    tbl2.save_pptx("test2.pptx")

    # data3 = [[1, 1], [0, 1], [2, 4], [3, 6]]
    tbl3 = PptxTable(data2)
    tbl3.set_formatting(first_row=True, first_col=False)
    tbl3.create_table(slide_index=0,
                      rows_sort_order=[2, 1, 0],
                      columns_sort_order=["pears", "bananas", "apples"],
                      columns_headers=["Pears", "Bananas", "Apples"],
                      columns_widths_weight=[1.7, .9, .9, .9],  # since transpose, we want to have 4 cols instead of 3
                      transpose=True)
    tbl3.save_pptx("test3.pptx")
